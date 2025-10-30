import streamlit as st
from docx import Document
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from collections import defaultdict
from copy import deepcopy
from io import BytesIO
from lxml import etree
import re
import unicodedata
from PIL import Image

# -------------------- CONFIGURAÇÃO INICIAL --------------------
st.set_page_config(layout="wide")
logo = Image.open("logo_jornada.png")
resample_filter = getattr(Image, "Resampling", Image).LANCZOS
logo = logo.resize((1235, 426), resample_filter)
_, col_logo, _ = st.columns([1, 4, 1])
with col_logo:
    st.image(logo, width=1235)
st.title("🚀 Gerador Automático de Slides")
st.info("CERTIFIQUE-SE DE ESTÁ FAZENDO O UPLOAD DOS ARQUIVOS CORRETOS ANTES DE GERAR OS SLIDES!")

# -------------------- FUNÇÕES AUXILIARES --------------------
def formatar_texto(texto, maiusculo_estado=False):
    texto = ' '.join(texto.strip().split())
    return texto.upper() if maiusculo_estado else ' '.join(w.capitalize() for w in texto.split())

def normalizar_texto_base(texto):
    if not texto:
        return ""
    texto = unicodedata.normalize("NFKD", str(texto))
    texto = "".join(ch for ch in texto if not unicodedata.combining(ch))
    texto = re.sub(r"\s+", " ", texto).strip()
    return texto.lower()

def sanitizar_nome_arquivo(nome):
    nome = (nome or "").strip()
    nome = re.sub(r'[\\/:*?"<>|]', "", nome)
    return nome or "Apresentacao_Final_Equipes"

def extrair_dados(uploaded_file):
    doc = Document(uploaded_file)
    registros = []
    for tabela in doc.tables:
        if not tabela.rows:
            continue

        cabecalho = [c.text.strip() for c in tabela.rows[0].cells]
        header_norm = [normalizar_texto_base(texto) for texto in cabecalho]

        aliases = {
            "Valido": [
                "valido",
                "alcance",
                "lancamentos validos",
                "alcance (m)",
                "distancia",
                "distancia (m)",
            ],
            "Equipe": [
                "equipe",
                "nome da equipe",
            ],
            "Funcao": [
                "funcao",
                "funcao/role",
                "funcao na equipe",
                "funcao integrante",
                "papel",
                "cargo",
            ],
            "Escola": [
                "escola",
                "nome da escola",
                "instituicao",
                "nome da instituicao",
                "colegio",
                "nome do colegio",
            ],
            "Cidade": [
                "cidade",
                "municipio",
            ],
            "Estado": [
                "estado",
                "uf",
            ],
            "Nome": [
                "nome",
                "nome do integrante",
                "nome integrante",
                "nome do aluno",
                "nome participante",
                "integrante",
                "participante",
                "aluno",
            ],
        }

        aliases_norm = {
            campo: [normalizar_texto_base(alias) for alias in lista]
            for campo, lista in aliases.items()
        }

        palavras_chave = {
            "Valido": {"alcance", "valido", "validos", "lancamento", "lancamentos", "distancia"},
            "Equipe": {"equipe", "time", "grupo"},
            "Funcao": {"funcao", "papel", "cargo"},
            "Escola": {"escola", "colegio", "instituicao"},
            "Cidade": {"cidade", "municipio"},
            "Estado": {"estado", "uf"},
            "Nome": {
                "nome",
                "nomes",
                "aluno",
                "alunos",
                "integrante",
                "integrantes",
                "participante",
                "participantes",
                "membro",
                "membros",
                "lider",
                "acompanhante",
                "responsavel",
                "responsaveis",
            },
        }

        tokens_por_coluna = []
        for cab_norm in header_norm:
            tokens = [tok for tok in re.split(r"[^a-z0-9]+", cab_norm) if tok]
            tokens_por_coluna.append(tokens)

        coluna_por_campo = {}
        colunas_usadas = set()

        def registrar(campo, idx):
            if idx is None or idx in colunas_usadas:
                return False
            coluna_por_campo[campo] = idx
            colunas_usadas.add(idx)
            return True

        # Correspondência exata com os aliases
        for campo, lista_aliases in aliases_norm.items():
            for alias_norm in lista_aliases:
                if not alias_norm:
                    continue
                for idx, cab_norm in enumerate(header_norm):
                    if idx in colunas_usadas:
                        continue
                    if cab_norm == alias_norm and registrar(campo, idx):
                        break
                if campo in coluna_por_campo:
                    break

        prioridade_campos = ["Valido", "Equipe", "Funcao", "Escola", "Cidade", "Estado", "Nome"]

        def combina(campo, tokens, cab_norm):
            if not cab_norm:
                return False
            tokens_set = set(tokens)
            chaves = palavras_chave.get(campo, set())
            if campo == "Nome":
                if tokens_set & {"escola", "colegio", "instituicao"}:
                    return False
            for chave in chaves:
                if chave in tokens_set:
                    return True
            for chave in chaves:
                if chave and chave in cab_norm:
                    return True
            return False

        for campo in prioridade_campos:
            if campo in coluna_por_campo:
                continue
            for idx, tokens in enumerate(tokens_por_coluna):
                if idx in colunas_usadas:
                    continue
                if combina(campo, tokens, header_norm[idx]):
                    registrar(campo, idx)
                    break

        def obter_valor(linha_celulas, chave):
            idx = coluna_por_campo.get(chave)
            if idx is not None and idx < len(linha_celulas):
                return linha_celulas[idx].strip()
            return ""

        for linha in tabela.rows[1:]:
            celulas = [c.text for c in linha.cells]
            if not any(c.strip() for c in celulas):
                continue

            registro = {chave: obter_valor(celulas, chave) for chave in aliases.keys()}

            if not registro["Equipe"] and not registro["Nome"]:
                continue

            registros.append({
                "Valido": registro["Valido"],
                "Equipe": registro["Equipe"],
                "Funcao": registro["Funcao"].lower(),
                "Escola": registro["Escola"],
                "Cidade": registro["Cidade"],
                "Estado": registro["Estado"],
                "Nome": registro["Nome"]
            })

    equipes = defaultdict(list)
    for r in registros:
        equipes[r["Equipe"]].append(r)

    def chave_ord(membros):
        try:
            return float(membros[0]["Valido"].replace(",", "."))
        except:
            return float("inf")

    equipes_ordenadas = sorted(equipes.items(), key=lambda x: chave_ord(x[1]))

    dados_finais = []
    for equipe_nome, membros in equipes_ordenadas:
        lider = [m for m in membros if "líder" in m["Funcao"] or "lider" in m["Funcao"]]
        acompanhante = [m for m in membros if "acompanhante" in m["Funcao"]]
        alunos = sorted(
            [m for m in membros if "aluno" in m["Funcao"]],
            key=lambda m: normalizar_texto_base(m["Nome"])
        )

        nomes_lider = formatar_texto(lider[0]["Nome"]) if lider else ""
        nomes_acompanhante = formatar_texto(acompanhante[0]["Nome"]) if acompanhante else ""

        linhas_nomes = []
        if nomes_lider:
            linhas_nomes.append(nomes_lider)
        if nomes_acompanhante:
            linhas_nomes.append(nomes_acompanhante)
        linhas_nomes += [formatar_texto(a["Nome"]) for a in alunos]

        nomes_formatados = "\n".join(linhas_nomes)

        info = membros[0]
        dados_finais.append({
            "{{LANCAMENTOS_VALIDOS}}": f"ALCANCE: {info['Valido']} m",
            "{{NOME_EQUIPE}}": f"Equipe: {equipe_nome.split()[-1]}",
            "{{NOME_ESCOLA}}": formatar_texto(info["Escola"]),
            "{{CIDADE_UF}}": f"{formatar_texto(info['Cidade'])} / {formatar_texto(info['Estado'], True)}",
            "{{NOMES_ALUNOS}}": nomes_formatados
        })
    return dados_finais

# -------------------- DUPLICAÇÃO DE SLIDE --------------------
def duplicate_slide_with_media(prs, source_slide):
    layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)
    for shape in source_slide.shapes:
        new_el = deepcopy(shape.element)
        if shape.shape_type == 13:  # picture
            try:
                img_blob = shape.image.blob
            except Exception:
                img_blob = None
            if img_blob:
                image_part, new_rId = new_slide.part.get_or_add_image_part(BytesIO(img_blob))
                new_el_xml = etree.fromstring(new_el.xml)
                blips = new_el_xml.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                for blip in blips:
                    blip.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', new_rId)
                from pptx.oxml import parse_xml
                new_el = parse_xml(etree.tostring(new_el_xml, encoding='utf-8'))
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    return new_slide

# -------------------- SUBSTITUIÇÃO DE PLACEHOLDERS --------------------
def replace_placeholders_in_shape(shape, team_data):
    if not shape.has_text_frame:
        return

    tf = shape.text_frame

    # Quando os placeholders de nomes e equipe estão na mesma caixa de texto, mas em
    # parágrafos diferentes, lidamos com todos de uma vez para garantir que as duas
    # informações sejam aplicadas com o mesmo estilo.
    frame_text = "\n".join("".join(run.text for run in paragraph.runs) for paragraph in tf.paragraphs)
    if "{{NOMES_ALUNOS}}" in frame_text and "{{NOME_EQUIPE}}" in frame_text:
        tf.clear()
        linhas = team_data["{{NOMES_ALUNOS}}"].split("\n") + [team_data["{{NOME_EQUIPE}}"]]
        for i, nome in enumerate(linhas):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            run = p.add_run()
            run.text = nome
            run.font.name = "Lexend"
            run.font.bold = True
            if i == len(linhas) - 1:  # última linha = nome da equipe
                run.font.size = Pt(20)
            else:
                run.font.size = Pt(26.5)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.CENTER
        return

    for paragraph in list(tf.paragraphs):
        full_text = "".join(run.text for run in paragraph.runs)

        # --- Corrige placeholders colados (ex: {{NOME_ESCOLA}}{{CIDADE_UF}} ou {{NOMES_ALUNOS}}{{NOME_EQUIPE}}) ---
        full_text = full_text.replace("}}{{", "}}\n{{")

        selected_key = None
        for k in team_data.keys():
            if k in full_text:
                selected_key = k
                break
        if not selected_key:
            continue

        # Substitui placeholders por valores
        new_text = full_text
        for k, v in team_data.items():
            new_text = new_text.replace(k, v)

        # Limpa runs antigos
        while paragraph.runs:
            paragraph._p.remove(paragraph.runs[0]._r)

        # --- ALCANCE ---
        if selected_key == "{{LANCAMENTOS_VALIDOS}}":
            match = re.match(r"(ALCANCE:\s*)([\d,.]+ m)", new_text, re.IGNORECASE)
            if match:
                prefix, valor = match.groups()
                run1 = paragraph.add_run()
                run1.text = prefix
                run1.font.name = "Lexend"
                run1.font.bold = False
                run1.font.size = Pt(28)
                run1.font.color.rgb = RGBColor(0x00, 0x6F, 0xC0)

                run2 = paragraph.add_run()
                run2.text = valor
                run2.font.name = "Lexend"
                run2.font.bold = True
                run2.font.underline = True
                run2.font.size = Pt(35)
                run2.font.color.rgb = RGBColor(0x00, 0x6F, 0xC0)

       # --- SOMENTE NOMES ---
        elif selected_key == "{{NOMES_ALUNOS}}":
            tf.clear()
            linhas = team_data["{{NOMES_ALUNOS}}"].split("\n")
            for i, nome in enumerate(linhas):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                run = p.add_run()
                run.text = nome
                run.font.name = "Lexend"
                run.font.bold = True
                run.font.size = Pt(26.5)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                p.alignment = PP_ALIGN.CENTER

        # --- NOME DA EQUIPE (se estiver sozinho) ---
        elif selected_key == "{{NOME_EQUIPE}}":
            run = paragraph.add_run()
            run.text = new_text
            run.font.name = "Lexend"
            run.font.bold = True
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            paragraph.alignment = PP_ALIGN.CENTER

        # --- ESCOLA + CIDADE ---
        elif "{{NOME_ESCOLA}}" in full_text and "{{CIDADE_UF}}" in full_text:
            tf.clear()
            partes = [team_data["{{NOME_ESCOLA}}"], team_data["{{CIDADE_UF}}"]]
            for i, parte in enumerate(partes):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                run = p.add_run()
                run.text = parte
                run.font.name = "Lexend"
                run.font.bold = True
                run.font.size = Pt(20)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                p.alignment = PP_ALIGN.CENTER

        # --- SOMENTE ESCOLA OU CIDADE (caso isolado) ---
        elif selected_key in ("{{NOME_ESCOLA}}", "{{CIDADE_UF}}"):
            run = paragraph.add_run()
            run.text = new_text
            run.font.name = "Lexend"
            run.font.bold = True
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            paragraph.alignment = PP_ALIGN.CENTER


# -------------------- GERAÇÃO FINAL --------------------
def gerar_apresentacao(dados, template_stream):
    prs = Presentation(template_stream)
    if not dados or not prs.slides:
        return prs

    modelo = prs.slides[0]
    slides_para_preencher = [modelo]

    for _ in range(len(dados) - 1):
        novo_slide = duplicate_slide_with_media(prs, modelo)
        slides_para_preencher.append(novo_slide)

    for slide, team in zip(slides_para_preencher, dados):
        for shape in slide.shapes:
            replace_placeholders_in_shape(shape, team)

    return prs

# -------------------- INTERFACE STREAMLIT --------------------
docx_file = st.file_uploader("📄 Arquivo DOCX", type=["docx"])
pptx_file = st.file_uploader("📊 Arquivo PPTX modelo", type=["pptx"])

if "nome_arquivo" not in st.session_state:
    st.session_state["nome_arquivo"] = ""

nome_arquivo_digitado = st.text_input(
    "Nome do arquivo (sem extensao)",
    value=st.session_state["nome_arquivo"],
)

if st.button("Confirmar nome do arquivo"):
    st.session_state["nome_arquivo"] = sanitizar_nome_arquivo(nome_arquivo_digitado)
    st.success(f"Nome para download ajustado para: {st.session_state['nome_arquivo']}.pptx")

if st.button("✨ Gerar Apresentação"):
    if not docx_file or not pptx_file:
        st.warning("Envie ambos os arquivos.")
    else:
        try:
            dados = extrair_dados(docx_file)
            if not dados:
                st.warning("Nenhum dado encontrado.")
            else:
                prs_final = gerar_apresentacao(dados, pptx_file)
                buf = BytesIO()
                prs_final.save(buf)
                buf.seek(0)
                st.success(f"Slides gerados: {len(dados)}")

                st.image("tiapamela.gif", caption="Apresentação pronta! 🚀", use_container_width=True)

                st.download_button(
                    "📥 Baixar Apresentação Final",
                    data=buf,
                    file_name=f"{st.session_state['nome_arquivo']}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )
        except Exception as e:
            st.error(f"Erro ao gerar apresentação: {e}")
